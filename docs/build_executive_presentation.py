import re
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


DOCS_DIR = Path(r"c:\MyGeneratedProjects\GitRepoPlan\repo-consolidated\docs")
SOURCE_FILES = [
    DOCS_DIR / "AI_Workflow_Blueprint.pptx",
    DOCS_DIR / "From_Snippets_to_Systems.pptx",
    DOCS_DIR / "The_Modular_AI_Workbench.pptx",
]
OUTPUT_FILE = DOCS_DIR / "Executive_AI_Context_Presentation.pptx"


def slide_number(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def slide_image_path(pptx_path: Path, slide_no: int) -> str | None:
    rels_name = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    with zipfile.ZipFile(pptx_path, "r") as zf:
        if rels_name not in zf.namelist():
            return None
        rels_xml = zf.read(rels_name).decode("utf-8", errors="ignore")
        image_match = re.search(r'Target="\.\./media/([^"]+)"', rels_xml)
        if not image_match:
            return None
        return f"ppt/media/{image_match.group(1)}"


def extract_image_bytes(pptx_path: Path, internal_image_path: str) -> bytes:
    with zipfile.ZipFile(pptx_path, "r") as zf:
        return zf.read(internal_image_path)


def pick_representative_slides(total_slides: int, pick_count: int = 3) -> list[int]:
    if total_slides <= pick_count:
        return list(range(1, total_slides + 1))
    picks = {1, total_slides}
    if pick_count >= 3:
        picks.add((total_slides + 1) // 2)
    if pick_count >= 4:
        picks.add(max(2, total_slides // 3))
    return sorted(picks)[:pick_count]


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)


def add_image_slide(prs: Presentation, img_path: Path, label: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    tx = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(6.8), Inches(0.4))
    tf = tx.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(10)


def build() -> None:
    temp_dir = DOCS_DIR / "_tmp_executive_images"
    temp_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Executive AI Context",
        "Condensed from three source decks for leadership review",
    )
    add_text_slide(
        prs,
        "Executive Agenda",
        [
            "Why this initiative and what changed",
            "Blueprint of the workflow and modular architecture",
            "Systemization journey and operating model",
            "Recommended next actions",
        ],
    )

    for src in SOURCE_FILES:
        if not src.exists():
            continue
        with zipfile.ZipFile(src, "r") as zf:
            slide_xmls = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        ordered = sorted(slide_xmls, key=slide_number)
        picks = pick_representative_slides(len(ordered), pick_count=3)

        add_text_slide(
            prs,
            src.stem.replace("_", " "),
            [
                f"Source deck has {len(ordered)} slides",
                "Selected representative slides: opening, core concept, closing",
                "See full-context deck for complete slide-by-slide coverage",
            ],
        )

        for slide_no in picks:
            img_internal = slide_image_path(src, slide_no)
            if not img_internal:
                continue
            img_bytes = extract_image_bytes(src, img_internal)
            img_file = temp_dir / f"{src.stem}_slide_{slide_no}{Path(img_internal).suffix}"
            img_file.write_bytes(img_bytes)
            add_image_slide(prs, img_file, f"{src.name} - executive selection (slide {slide_no})")

    add_text_slide(
        prs,
        "Recommended Next Actions",
        [
            "Standardize reusable components into a shared module catalog",
            "Define governance for prompts, agents, and workflow quality gates",
            "Measure adoption and impact with release and productivity metrics",
            "Run quarterly architecture and model strategy review",
        ],
    )

    prs.save(str(OUTPUT_FILE))


if __name__ == "__main__":
    build()
