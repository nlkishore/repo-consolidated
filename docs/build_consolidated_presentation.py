import re
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


DOCS_DIR = Path(r"c:\MyGeneratedProjects\GitRepoPlan\repo-consolidated\docs")
SOURCE_FILES = [
    DOCS_DIR / "AI_Workflow_Blueprint.pptx",
    DOCS_DIR / "From_Snippets_to_Systems.pptx",
    DOCS_DIR / "The_Modular_AI_Workbench.pptx",
]
OUTPUT_FILE = DOCS_DIR / "Consolidated_AI_Context_Presentation.pptx"


def slide_number(name: str) -> int:
    match = re.search(r"slide(\d+)\.xml$", name)
    return int(match.group(1)) if match else 0


def slide_image_path(pptx_path: Path, slide_no: int) -> str | None:
    rels_name = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    if not rels_name:
        return None
    with zipfile.ZipFile(pptx_path, "r") as zf:
        if rels_name not in zf.namelist():
            return None
        rels_xml = zf.read(rels_name).decode("utf-8", errors="ignore")
        image_match = re.search(r'Target="\.\./media/([^"]+)"', rels_xml)
        if not image_match:
            return None
        return f"ppt/media/{image_match.group(1)}"


def extract_image_bytes(pptx_path: Path, internal_image_path: str) -> bytes:
    with zipfile.ZipFile(pptx_path, "r") as zf:
        return zf.read(internal_image_path)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.0), Inches(1.2))
    tf = tx.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(24)


def add_image_slide(prs: Presentation, img_path: Path, label: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    tag = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(6.5), Inches(0.4))
    tf = tag.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(10)


def build() -> None:
    temp_dir = DOCS_DIR / "_tmp_consolidated_images"
    temp_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Consolidated AI Context",
        "Merged from: AI Workflow Blueprint, From Snippets to Systems, The Modular AI Workbench",
    )
    add_section_slide(
        prs,
        "How to Use This Deck",
        "This consolidated deck preserves every source slide in sequence with source labels for full traceability.",
    )

    for src in SOURCE_FILES:
        if not src.exists():
            continue
        with zipfile.ZipFile(src, "r") as zf:
            slide_xmls = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        ordered = sorted(slide_xmls, key=slide_number)
        add_section_slide(prs, src.stem.replace("_", " "), f"{len(ordered)} slides imported")

        for slide_xml in ordered:
            s_no = slide_number(slide_xml)
            img_internal = slide_image_path(src, s_no)
            if not img_internal:
                continue
            img_bytes = extract_image_bytes(src, img_internal)
            img_suffix = Path(img_internal).suffix or ".png"
            img_file = temp_dir / f"{src.stem}_slide_{s_no}{img_suffix}"
            img_file.write_bytes(img_bytes)
            add_image_slide(prs, img_file, f"{src.name} - slide {s_no}")

    prs.save(str(OUTPUT_FILE))


if __name__ == "__main__":
    build()
