from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


OUT_PATH = Path(
    r"c:\MyGeneratedProjects\GitRepoPlan\repo-consolidated\docs\Team_AI_Tools_Evolution_Awareness_Session.pptx"
)


def add_title(prs: Presentation, title: str, subtitle: str, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.notes_slide.notes_text_frame.text = notes


def add_bullets(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(22)
    slide.notes_slide.notes_text_frame.text = notes


def add_two_column(prs: Presentation, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.9), Inches(5.4))
    ltf = left_box.text_frame
    ltf.clear()
    p0 = ltf.paragraphs[0]
    p0.text = left_title
    p0.font.bold = True
    p0.font.size = Pt(24)
    for item in left_items:
        p = ltf.add_paragraph()
        p.text = f"- {item}"
        p.level = 1
        p.font.size = Pt(20)

    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.4))
    rtf = right_box.text_frame
    rtf.clear()
    q0 = rtf.paragraphs[0]
    q0.text = right_title
    q0.font.bold = True
    q0.font.size = Pt(24)
    for item in right_items:
        q = rtf.add_paragraph()
        q.text = f"- {item}"
        q.level = 1
        q.font.size = Pt(20)

    slide.notes_slide.notes_text_frame.text = notes


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title(
        prs,
        "AI Tools Evolution Awareness Session",
        "Preparing our team for organizational adoption of advanced AI tools",
        "Welcome everyone. Today is about understanding where we are now, where AI development tools are heading, and how we prepare as a team."
    )

    add_bullets(
        prs,
        "Session Objectives",
        [
            "Understand 3 generations of AI coding tools",
            "Map current daily work use cases to each generation",
            "Identify capability gaps in our team today",
            "Define a practical upskilling and adoption plan",
        ],
        "Set expectations: this is a practical session. We are not just discussing trends; we are defining concrete team actions."
    )

    add_bullets(
        prs,
        "Current State at Work (Today)",
        [
            "Primary tool: GitHub Copilot workflow",
            "Current limit: up to 3 files per review/summarization cycle",
            "Manual steps needed: copy suggestions into project",
            "Partial acceleration, but high manual integration effort",
        ],
        "Acknowledge value of current tooling, but highlight friction points: context limits and manual copy-paste reduce end-to-end productivity."
    )

    add_two_column(
        prs,
        "Generation 1: Assistive Copilot",
        "What it does well",
        [
            "Inline suggestions and quick summaries",
            "Helps with local refactoring ideas",
            "Speeds up repetitive coding tasks",
        ],
        "Where it slows us down",
        [
            "Limited project context window",
            "Manual patching into multiple files",
            "No autonomous validation + fix loop",
        ],
        "Explain that Gen1 is assistant mode. Human still orchestrates everything: analysis, implementation, testing, and final wiring."
    )

    add_bullets(
        prs,
        "Generation 2: Prompt-to-Outcome Tools",
        [
            "Input: use case or prompt with acceptance criteria",
            "Output: code + tests + bug fixes + validated result",
            "Tool can iterate until tests pass",
            "Team focus shifts from typing code to defining intent",
        ],
        "This is a major shift: from code completion to task completion. Emphasize the need for better prompt engineering and review discipline."
    )

    add_two_column(
        prs,
        "Generation 2 in Daily Team Work",
        "High-value use cases",
        [
            "Bug triage and root-cause fixes",
            "Add endpoint + test coverage",
            "Refactor with regression safeguards",
            "Test generation for legacy modules",
        ],
        "Guardrails needed",
        [
            "Clear Definition of Done",
            "Automated test gates in CI",
            "Security/static checks before merge",
            "Human review for architecture fit",
        ],
        "Clarify that Gen2 works best with explicit constraints and quality gates. Without guardrails, speed can create risk."
    )

    add_bullets(
        prs,
        "Generation 3: Agentic AI (Multi-Use-Case Orchestration)",
        [
            "Agents combine multiple tasks/features into one flow",
            "Can plan, execute, test, and integrate across modules",
            "Supports end-to-end project scaffolding and evolution",
            "Requires stronger governance and architecture standards",
        ],
        "Frame Gen3 as system-level automation. Team role becomes orchestration, policy setting, and outcome assurance."
    )

    add_two_column(
        prs,
        "Skill Shift Required for Our Team",
        "From",
        [
            "Manual coding-first execution",
            "Single-task tooling mindset",
            "Ad hoc review practices",
        ],
        "To",
        [
            "Intent-first specification writing",
            "Agent/task orchestration skills",
            "Quality gates + policy-driven workflows",
        ],
        "This is a capability transition. We need to train for prompting, evaluation, and workflow design, not only language syntax."
    )

    add_bullets(
        prs,
        "90-Day Enablement Roadmap",
        [
            "Weeks 1-3: Baseline training on prompt and review patterns",
            "Weeks 4-6: Pilot Gen2 workflows on 2-3 real use cases",
            "Weeks 7-9: Measure throughput, defects, and cycle time",
            "Weeks 10-12: Introduce Gen3 agent pilots with governance",
        ],
        "Position this as phased adoption. Start with measurable pilots, then scale based on evidence."
    )

    add_two_column(
        prs,
        "Team Use Case Mapping (Example)",
        "Current recurring work",
        [
            "Defect fixes in existing modules",
            "New feature endpoint development",
            "Refactoring + test gap closure",
            "Release readiness verification",
        ],
        "Best-fit AI generation",
        [
            "Gen1/Gen2 depending on scope",
            "Gen2 for faster code+test delivery",
            "Gen2 with strict regression checks",
            "Gen3 agents for release orchestration",
        ],
        "Use this as discussion slide. Ask team to map their own backlog items to Gen1/2/3 after the session."
    )

    add_two_column(
        prs,
        "Free AI Tools by Role",
        "Developers",
        [
            "Cursor Free tier / VS Code AI extensions",
            "Continue.dev (open-source coding assistant)",
            "Codeium free plan",
            "Aider (CLI pair-programming with local/remote models)",
        ],
        "System Analysts + Business Analysts",
        [
            "ChatGPT free / Gemini free for requirement drafting",
            "Notion AI free trial workflows for documentation",
            "Perplexity free for research with citations",
            "Miro AI / FigJam AI basic features for workshops",
        ],
        "Clarify that tool availability can change. Team should validate current licensing, data policy, and enterprise restrictions before adoption."
    )

    add_two_column(
        prs,
        "System Analyst Use Cases (Practical)",
        "Analysis and design",
        [
            "Convert BRD to functional requirements",
            "Generate API contract drafts from use cases",
            "Create sequence diagrams from textual flows",
            "Summarize impact analysis across modules",
        ],
        "Delivery acceleration",
        [
            "Prepare user-story acceptance criteria",
            "Draft test scenarios and edge cases",
            "Generate release note skeletons",
            "Create dependency and risk checklists",
        ],
        "For System Analysts, focus on quality of requirements, traceability, and faster handover to dev and QA."
    )

    add_two_column(
        prs,
        "Business Analyst Use Cases (Practical)",
        "Business discovery",
        [
            "Interview summary into structured requirements",
            "Persona and journey map first drafts",
            "Process pain-point extraction from call notes",
            "Stakeholder communication draft generation",
        ],
        "Decision support",
        [
            "Option comparison tables with pros/cons",
            "Benefit hypothesis and KPI draft baselines",
            "Backlog prioritization rationale draft",
            "Business-case narrative for leadership updates",
        ],
        "For Business Analysts, AI helps reduce documentation latency and improve consistency in stakeholder communication."
    )

    add_bullets(
        prs,
        "Safe Adoption Checklist for Free Tools",
        [
            "Never paste production secrets or PII into public tools",
            "Use anonymized samples for requirement or defect analysis",
            "Keep final business decisions human-reviewed",
            "Store prompt templates and approved patterns in team wiki",
        ],
        "Position this as mandatory governance. Free tools are useful for learning and prototyping, but policy compliance is non-negotiable."
    )

    add_bullets(
        prs,
        "Risks and Controls",
        [
            "Risk: over-reliance on generated code without review",
            "Control: mandatory architecture and security checkpoints",
            "Risk: inconsistent prompts and outcomes",
            "Control: reusable prompt templates and playbooks",
        ],
        "Reassure the team that controls are part of scaling safely, not slowing down innovation."
    )

    add_bullets(
        prs,
        "Call to Action",
        [
            "Nominate pilot owners for top 3 use cases",
            "Define success metrics per use case this week",
            "Start Gen2 pilot in next sprint",
            "Review Gen3 readiness at end of quarter",
        ],
        "Close with ownership and timeline. Ask for volunteers and alignment before ending the session."
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build()
